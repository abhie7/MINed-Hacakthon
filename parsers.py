from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import re


class PdfParser:
    '''Extract text from a PDF file.'''
    def __init__(self, filepath: str):
        self.filepath = filepath

    def extract_text(self) -> str:
        reader = PdfReader(self.filepath)  # create a PdfReader object
        text = '\n'.join(page.extract_text() for page in reader.pages)  # extract text from each page
        return text # return the extracted text

class PptxParser:
    '''Extract text from a PPT.'''
    def __init__(self, filepath: str):
        self.filepath = filepath

    def extract_text(self) -> str:
        presentation = Presentation(self.filepath) # creating a Presentation object
        # extract text from each shape in each slide
        # that has a text frame, then join the results with '\n' as the separator
        text = '\n'.join(
            shape.text_frame.text
            for slide in presentation.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        return text  # return the extracted text

class DocxParser:
    '''Extract text from a .docx file'''
    def __init__(self, filepath: str):
        self.filepath = filepath

    def extract_text(self) -> str:
        document = Document(self.filepath) # create a Document object
        # uses list comprehension to collect all paragraph texts
        text = '\n'.join(paragraph.text for paragraph in document.paragraphs)
        return text # return the extracted text

# Comment this module if it shows an error - not tested
class TexParser:
    '''Extract text from a .tex file'''
    def __init__(self, filepath: str):
        self.filepath = filepath

    def extract_text(self) -> str:
        with open(self.filepath, 'r', encoding='utf-8') as tex_file:
            tex_content = tex_file.read()

        text = re.sub(r'\\[^{}]+({[^{}]})', '', tex_content)
        text = re.sub(r'%.*', '', text)
        text = re.sub(r'\s+', ' ', text)
        text = text.strip()

        return text

# test
if __name__ == "__main__":
    extractor = TexParser("./data/files/sample.tex")
    extracted_text = extractor.extract_text()
    print(extracted_text)